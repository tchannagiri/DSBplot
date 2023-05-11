import pptx
import pptx.util
import pptx.chart.data
import pptx.enum.shapes
import pptx.dml.color
import pptx.dml.effect

import pptx.shapes.connector
import pptx.enum.shapes
import pptx.enum.text
import pptx.enum.dml
import pptx.dml.color
import pptx.util


def add_textbox_pptx(
  slide,
  text,
  x_pt,
  y_pt,
  width_pt,
  height_pt,
  font_size_pt,
  text_align = 'center',
):
  shape = slide.shapes.add_shape(
    pptx.enum.shapes.MSO_SHAPE.RECTANGLE,
    pptx.util.Pt(x_pt),
    pptx.util.Pt(y_pt),
    pptx.util.Pt(width_pt),
    pptx.util.Pt(height_pt),
  )
  paragraph = shape.text_frame.paragraphs[0]
  paragraph.alignment = {
    'left': pptx.enum.text.PP_PARAGRAPH_ALIGNMENT.LEFT,
    'right': pptx.enum.text.PP_PARAGRAPH_ALIGNMENT.RIGHT,
    'center': pptx.enum.text.PP_PARAGRAPH_ALIGNMENT.CENTER,
  }[text_align]
  run = paragraph.add_run()
  run.text = text
  run.font.size = pptx.util.Pt(font_size_pt)
  run.font.color.rgb = pptx.dml.color.RGBColor(0, 0, 0)
  shape.fill.background()
  shape.line.fill.background()
  return shape

def add_textbox_grid_pptx(
  slide,
  text_grid,
  x_pt,
  y_pt,
  cell_width_pt,
  cell_height_pt,
  cell_width_spacing_pt,
  cell_height_spacing_pt,
  text_width_pt,
  text_height_pt,
  font_size_pt,
  text_align = 'center',
):
  for row in range(text_grid.shape[0]):
    for col in range(text_grid.shape[1]):
      x_pt_cell = x_pt + col * (cell_width_pt + cell_width_spacing_pt)
      y_pt_cell = y_pt + row * (cell_height_pt + cell_height_spacing_pt)
      add_textbox_pptx(
        slide = slide,
        text = text_grid[row, col],
        x_pt = x_pt_cell,
        y_pt = y_pt_cell,
        width_pt = text_width_pt,
        height_pt = text_height_pt,
        font_size_pt = font_size_pt,
        text_align = text_align,
      )

def add_picture_pptx(
  slide,
  file_name,
  x_pt,
  y_pt,
  width_pt,
  height_pt,
):
  slide.shapes.add_picture(
    file_name,
    pptx.util.Pt(x_pt),
    pptx.util.Pt(y_pt),
    pptx.util.Pt(width_pt),
    pptx.util.Pt(height_pt),
  )

def add_picture_grid_pptx(
  slide,
  file_name_grid,
  x_pt,
  y_pt,
  cell_width_pt,
  cell_height_pt,
  cell_width_spacing_pt,
  cell_height_spacing_pt,
):
  for row in range(file_name_grid.shape[0]):
    for col in range(file_name_grid.shape[1]):
      add_picture_pptx(
        slide = slide,
        file_name = file_name_grid[row, col],
        x_pt = x_pt + col * (cell_width_pt + cell_width_spacing_pt),
        y_pt = y_pt + row * (cell_height_pt + cell_height_spacing_pt),
        width_pt = cell_width_pt,
        height_pt = cell_height_pt,
      )