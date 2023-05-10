import os
import sys
sys.path.append(os.path.abspath(os.path.join(os.path.dirname(__file__), '../utils/'))) # allow importing the utils dir
sys.path.append(os.path.abspath(os.path.join(os.path.dirname(__file__), '../2_graph_processing/'))) # allow importing the graphs dir
import argparse

import numpy as np

import common_utils
import file_utils
import log_utils
import get_pptx_helpers

import pptx
import pptx.util
import pptx.enum.shapes
import pptx.enum.text
import pptx.dml.color
import pptx.dml.effect
import PIL

PPTX_TEMPLATE_FILE = os.path.join(os.path.dirname(__file__), 'template.pptx') # make this an arg!

### Font sizes ###
IMAGE_LABEL_FONT_SIZE_PT = 8
TITLE_FONT_SIZE_PT = 18
IMAGE_LABEL_WIDTH_PT = 60
IMAGE_LABEL_HEIGHT_PT = 20
MARGIN_CORNER_LABEL_FONT_SIZE_PT = 12
MARGIN_TOP_LABEL_FONT_SIZE_PT = 9
MARGIN_LEFT_LABEL_FONT_SIZE_PT = 8
LEGEND_TITLE_FONT_SIZE_PT = 10
LEGEND_LABEL_FONT_SIZE_PT = 8

### Height/width of elements ###
TITLE_HEIGHT_PT = 30
MARGIN_TOP_HEIGHT_PT = 20
MARGIN_LEFT_WIDTH_PT = 60
MARGIN_LEFT_SPILL_OVER_PT = 10
MARGIN_RIGHT_SPILL_OVER_PT = 10
MARGIN_RIGHT_WIDTH_PT = 100
IMAGE_HEIGHT_SPACING_PT = 10
IMAGE_WIDTH_SPACING_PT = 5
GRID_HEIGHT_SPACING_PT = 40

def get_slide(
  prs,
  image_grid_list,
  total_width_frac_list = None,
  title = None,
  title_height_pt = TITLE_HEIGHT_PT,
  title_font_size_pt = TITLE_FONT_SIZE_PT,
  image_label_grid_list = None,
  image_label_font_size_pt = IMAGE_LABEL_FONT_SIZE_PT,
  image_label_width_pt = IMAGE_LABEL_WIDTH_PT,
  image_label_height_pt = IMAGE_LABEL_HEIGHT_PT,
  margin_label_top_list = None,
  margin_label_left_list = None,
  margin_top_font_size_pt = MARGIN_TOP_LABEL_FONT_SIZE_PT,
  margin_left_font_size_pt = MARGIN_LEFT_LABEL_FONT_SIZE_PT,
  margin_top_height_pt = MARGIN_TOP_HEIGHT_PT,
  margin_left_width_pt = MARGIN_LEFT_WIDTH_PT,
  margin_left_spill_over_pt = MARGIN_LEFT_SPILL_OVER_PT,
  margin_right_spill_over_pt = MARGIN_RIGHT_SPILL_OVER_PT,
  image_height_spacing_pt = IMAGE_HEIGHT_SPACING_PT,
  image_width_spacing_pt = IMAGE_WIDTH_SPACING_PT,
  grid_height_spacing_pt = GRID_HEIGHT_SPACING_PT,
):
  num_grids = len(image_grid_list)

  if total_width_frac_list is None:
    total_width_frac_list = [1] * num_grids

  image_label_show = image_label_grid_list is not None
  if image_label_show and (len(image_label_grid_list) != num_grids):
    raise Exception(
      f"Incorrect number of image label grids : {len(image_label_grid_list)}." +
      f" Expected {num_grids}."
    )

  margin_show_top = margin_label_top_list is not None
  margin_show_left = margin_label_left_list is not None
  if margin_show_left:
    if len(margin_label_left_list) != num_grids:
      raise Exception(
        f"Incorrect number of left margin lists: {len(margin_label_left_list)}." +
        f" Expected {num_grids}."
      )
  if margin_show_top:
    if len(margin_label_top_list) != num_grids:
      raise Exception(
        f"Incorrect number of top margin lists: {len(margin_label_top_list)}." +
        f" Expected {num_grids}."
      )

  title_slide_layout = prs.slide_layouts[6] # should be a blank layout
  slide = prs.slides.add_slide(title_slide_layout)
  
  slide_width_pt = prs.slide_width / pptx.util.Pt(1)
  slide_height_pt = prs.slide_height / pptx.util.Pt(1)

  log_utils.log(f"Slide width: {slide_width_pt} pt")
  log_utils.log(f"Slide height: {slide_height_pt} pt")

  y_pt = 0
  x_pt = 0

  # Title
  if title is not None:
    get_pptx_helpers.add_textbox_pptx(
      slide = slide,
      text = title,
      x_pt = 0,
      y_pt = 0,
      width_pt = slide_width_pt,
      height_pt = title_height_pt,
      font_size_pt =  title_font_size_pt,
    )
    y_pt += title_height_pt

  for i in range(num_grids):
    max_image_width_px = 0
    max_image_height_px = 0

    # Get point/pixel ratio and dimensions
    # All images should be the same width/height
    for row in range(image_grid_list[i].shape[0]):
      for col in range(image_grid_list[i].shape[1]):
        image = PIL.Image.open(image_grid_list[i][row][col])
        max_image_width_px = max(max_image_width_px, image.width)
        max_image_height_px = max(max_image_height_px, image.height)

    content_width_px = image_grid_list[i].shape[1] * max_image_width_px

    content_width_pt = slide_width_pt * total_width_frac_list[i]
    if margin_show_left:
      content_width_pt -= margin_left_width_pt

    ratio_pt_px = content_width_pt / content_width_px
    cell_width_pt = ratio_pt_px * max_image_width_px
    cell_height_pt = ratio_pt_px * max_image_height_px

    content_x_pt = x_pt
    content_y_pt = y_pt
    content_height_pt = (
      (image_grid_list[i].shape[0] * cell_height_pt) + 
      ((image_grid_list[i].shape[0] - 1) * image_height_spacing_pt)
    )
    content_width_pt = (
      (image_grid_list[i].shape[1] * cell_width_pt) + 
      ((image_grid_list[i].shape[1] - 1) * image_width_spacing_pt)
    )

    if margin_show_top:
      content_y_pt += margin_top_height_pt
    if margin_show_left:
      content_x_pt += margin_left_width_pt

    # Content images
    get_pptx_helpers.add_picture_grid_pptx(
      slide = slide,
      file_name_grid = image_grid_list[i],
      x_pt = content_x_pt,
      y_pt = content_y_pt,
      cell_width_pt = cell_width_pt,
      cell_height_pt = cell_height_pt,
      cell_width_spacing_pt = image_width_spacing_pt,
      cell_height_spacing_pt = image_height_spacing_pt,
    )

    if image_label_show:
      if image_grid_list[i].shape != image_grid_list[i].shape:
        raise Exception(
          f"Incorrect shape of {i}'th label grid: {image_grid_list[i].shape}." +
          f" Expected: {image_grid_list[i].shape}."
        )
      # Image labels
      get_pptx_helpers.add_textbox_grid_pptx(
        slide = slide,
        text_grid = image_label_grid_list[i],
        x_pt = content_x_pt,
        y_pt = content_y_pt,
        cell_width_pt = cell_width_pt,
        cell_height_pt = cell_height_pt,
        cell_width_spacing_pt = image_width_spacing_pt,
        cell_height_spacing_pt = image_height_spacing_pt,
        text_width_pt = image_label_width_pt,
        text_height_pt = image_label_height_pt,
        font_size_pt = image_label_font_size_pt,
        text_align = 'center',
      )

    if margin_show_top:
      # Top margin
      get_pptx_helpers.add_textbox_grid_pptx(
        slide = slide,
        text_grid = np.array(margin_label_top_list[i]).reshape(1, -1),
        x_pt = content_x_pt,
        y_pt = y_pt,
        cell_width_pt = cell_width_pt,
        cell_height_pt = margin_top_height_pt,
        cell_width_spacing_pt = image_width_spacing_pt,
        cell_height_spacing_pt = image_height_spacing_pt,
        text_height_pt = margin_top_height_pt,
        text_width_pt = cell_width_pt,
        font_size_pt = margin_top_font_size_pt,
      )

    if margin_show_left:
      # Left margin
      cell_width_pt_left_margin = (
        margin_left_width_pt +
        margin_left_spill_over_pt +
        margin_right_spill_over_pt
      )
      get_pptx_helpers.add_textbox_grid_pptx(
        slide = slide,
        text_grid = np.array(margin_label_left_list[i]).reshape(-1, 1),
        x_pt = -margin_left_spill_over_pt,
        y_pt = content_y_pt,
        cell_width_pt = cell_width_pt_left_margin,
        cell_height_pt = cell_height_pt,
        cell_width_spacing_pt = image_width_spacing_pt,
        cell_height_spacing_pt = image_height_spacing_pt,
        text_height_pt = cell_height_pt,
        text_width_pt = cell_width_pt_left_margin,
        font_size_pt = margin_left_font_size_pt,
      )

    y_pt = content_y_pt + content_height_pt + grid_height_spacing_pt
  # End grid for loop

def parse_args():
  parser = argparse.ArgumentParser(
    description = 'Create powerpoint figures from images.',
    formatter_class = argparse.ArgumentDefaultsHelpFormatter,
  )
  parser.add_argument(
    '--input',
    nargs = '+',
    type = common_utils.check_file,
    help = (
      'List of images to include in the grid.' +
      ' The number of arguments should match the number of grids' +
      ' and their dimensions.'
    ),
    required = True,
  )
  parser.add_argument(
    '--num_grids',
    required = True,
    type = int,
    help = 'Number of separate grids to create.',
  )
  parser.add_argument(
    '--num_rows',
    nargs = '+',
    required = True,
    type = int,
    help = (
      'Number of rows in each grid.' +
      ' Number of arguments should match the number of grids.'
    ),
  )
  parser.add_argument(
    '--num_cols',
    nargs = '+',
    required = True,
    type = int,
    help = (
      'Number of columns in each grid.' +
      ' Number of arguments should match the number of grids.'
    ),
  )
  parser.add_argument(
    '--image_height_spacing_pt',
    type = float,
    default = IMAGE_HEIGHT_SPACING_PT,
    help = (
      'How much vertical spacing between rows of images' +
      ' in the same grid in pt.'
    ),
  )
  parser.add_argument(
    '--image_width_spacing_pt',
    type = float,
    default = IMAGE_WIDTH_SPACING_PT,
    help = (
      'How much horizontal spacing between columns of images' +
      ' in the same grid in pt.'
    ),
  )
  parser.add_argument(
    '--grid_height_spacing_pt',
    type = float,
    default = GRID_HEIGHT_SPACING_PT,
    help = 'How much vertical spacing between different grids in pt.',
  )
  parser.add_argument(
    '--image_labels',
    nargs = '+',
    help = (
      'Labels of images in the grid.' +
      ' The number of arguments should match the number of grids' +
      ' and their dimensions.' +
      ' The label textboxs will be positioned so that the top left corner' +
      ' of the textbox is at the top left corner of the image.'
    ),
  )
  parser.add_argument(
    '--image_label_width_pt',
    type = float,
    default = IMAGE_LABEL_WIDTH_PT,
    help = 'Width of image label box in pt.',
  )
  parser.add_argument(
    '--image_label_height_pt',
    type = float,
    default = IMAGE_LABEL_HEIGHT_PT,
    help = 'Height of image label box in pt.',
  )
  parser.add_argument(
    '--image_label_font_size_pt',
    type = float,
    default = IMAGE_LABEL_FONT_SIZE_PT,
    help = 'Size of image label font in pt.',
  )
  parser.add_argument(
    '--margin_labels_top',
    nargs = '+',
    help = (
      'Labels in the top margins of the grids.' +
      ' The number of arguments must be sum of the number of columns in each grid.'
    ),
  )
  parser.add_argument(
    '--margin_top_height_pt',
    type = float,
    default = MARGIN_TOP_HEIGHT_PT,
    help = (
      'Height of top margin in pt.' +
      ' The top margin is where the top labels are placed.'
    ),
  )
  parser.add_argument(
    '--margin_top_font_size_pt',
    type = float,
    default = MARGIN_TOP_LABEL_FONT_SIZE_PT,
    help = 'Size of top margin label font in pt.',
  )
  parser.add_argument(
    '--margin_labels_left',
    nargs = '+',
    help = (
      'Labels in the left margins of the grids,' +
      ' The number of arguments must be sum of the number of rows in each grid.'
    ),
  )
  parser.add_argument(
    '--margin_left_width_pt',
    type = float,
    default = MARGIN_LEFT_WIDTH_PT,
    help = (
      'Width of left margin in pt.' +
      ' The left margin is where the left labels are placed.'
    ),
  )
  parser.add_argument(
    '--margin_left_spill_over_pt',
    type = float,
    nargs = 2,
    default = MARGIN_LEFT_SPILL_OVER_PT,
    help = (
      'How much to extend the left label boxes on the left and right in pt.' +
      ' This is to help centering the labels.'
    ),
  )
  parser.add_argument(
    '--margin_left_font_size_pt',
    type = float,
    default = MARGIN_LEFT_LABEL_FONT_SIZE_PT,
    help = 'Size of left margin label font in pt.',
  )
  parser.add_argument(
    '--total_width_frac',
    nargs = '+',
    type = float,
    help = (
      'Fraction of the page width to use for each grid .' +
      ' Number of arguments should match the number of grids.'
    ),
  )
  parser.add_argument(
    '--title',
    help = 'Page title.',
  )
  parser.add_argument(
    '--title_height_pt',
    type = float,
    default = TITLE_HEIGHT_PT,
    help = (
      'Height of title box in pt.' +
      ' The title box spans the whole page width.'
    ),
  )
  parser.add_argument(
    '--title_font_size_pt',
    type = float,
    default = TITLE_FONT_SIZE_PT,
    help = 'Size of title font in pt.',
  )
  parser.add_argument(
    '--template',
    default = os.path.join(os.path.dirname(__file__), 'template.pptx'),
    help = 'The PPTX file to use as a template. This determines the page size.',
  )
  parser.add_argument(
    '--output',
    help = 'Output PPTX file.',
    required = True,
  )
  return vars(parser.parse_args())

def main(
  input,
  num_grids,
  num_rows,
  num_cols,
  image_height_spacing_pt,
  image_width_spacing_pt,
  grid_height_spacing_pt,
  image_labels,
  image_label_width_pt,
  image_label_height_pt,
  image_label_font_size_pt,
  margin_labels_top,
  margin_top_height_pt,
  margin_top_font_size_pt,
  margin_labels_left,
  margin_left_width_pt,
  margin_left_spill_over_pt,
  margin_left_font_size_pt,
  total_width_frac,
  title,
  title_height_pt,
  title_font_size_pt,
  template,
  output,
):
  if num_grids != len(num_rows):
    raise Exception(
      f'Incorrect num rows specification: {num_rows}.' +
      f' Expected {num_grids} values.' 
    )
  
  if num_grids != len(num_cols):
    raise Exception(
      f'Incorrect num columns specification: {num_cols}.' +
      f' Expected {num_grids} values.' 
    )

  if total_width_frac is None:
    total_width_frac = [1] * num_grids
  if len(total_width_frac) != num_grids:
    raise Exception(
      f'Incorrect number of total widths: {total_width_frac}.' +
      f' Expected {num_grids} values.' 
    )

  num_images_total = sum(r * c for r, c in zip(num_rows, num_cols))
  if num_images_total != len(input):
    raise Exception(
      f'Incorrect number of input files: {len(input)}.' +
      f' Expected {num_images_total} values.' 
    )

  image_grid_list = []
  image_index = 0
  for i in range(num_grids):
    num_images = num_rows[i] * num_cols[i]
    image_grid = np.array(
      [input[image_index + j] for j in range(num_images)]
    )
    image_grid = image_grid.reshape((num_rows[i], num_cols[i]))
    image_grid_list.append(image_grid)
    image_index += num_images

  label_grid_list = None
  if image_labels is not None:
    if len(image_labels) != num_images_total:
      raise Exception(
        f'Incorrect number of input labels: {len(input)}.' +
        f' Expected {num_images_total} values.' 
      )
    label_grid_list = []
    index = 0
    for i in range(num_grids):
      num_labels = num_rows[i] * num_cols[i]
      label_grid = np.array(
        [image_labels[index + j] for j in range(num_labels)]
      )
      label_grid = label_grid.reshape((num_rows[i], num_cols[i]))
      label_grid_list.append(label_grid)
      index += num_labels
  
  margin_label_left_list = None
  if margin_labels_left is not None:
    num_margin_labels_left = sum(num_rows[i] for i in range(num_grids))
    if len(margin_labels_left) != num_margin_labels_left:
      raise Exception(
        f'Incorrect number of left margin labels: {len(margin_labels_left)}.' +
        f' Expected {num_margin_labels_left} values.' 
      )
    margin_label_left_list = []
    index = 0
    for i in range(num_grids):
      num_labels = num_rows[i]
      label_list = [margin_labels_left[index + j] for j in range(num_labels)]
      margin_label_left_list.append(label_list)
      index += num_labels
  
  margin_label_top_list = None
  if margin_labels_top is not None:
    num_margin_labels_top = sum(num_cols[i] for i in range(num_grids))
    if len(margin_labels_top) != num_margin_labels_top:
      raise Exception(
        f'Incorrect number of top margin labels: {len(margin_labels_top)}.' +
        f' Expected {num_margin_labels_top} values.' 
      )
    margin_label_top_list = []
    index = 0
    for i in range(num_grids):
      num_labels = num_cols[i]
      label_list = [margin_labels_top[index + j] for j in range(num_labels)]
      margin_label_top_list.append(label_list)
      index += num_labels

  # replace \n with newline characters in case using Windows
  for i in range(num_grids):
    if label_grid_list is not None:
      for r in range(num_rows[i]):
        for c in range(num_cols[i]):
          label_grid_list[i][r, c] = label_grid_list[i][r, c].replace('\\n', '\n')
    if margin_label_left_list is not None:
      for r in range(num_rows[i]):
        margin_label_left_list[i][r] = margin_label_left_list[i][r].replace('\\n', '\n')
    if margin_label_top_list is not None:
      for c in range(num_cols[i]):
        margin_label_top_list[i][c] = margin_label_top_list[i][c].replace('\\n', '\n')

  prs = pptx.Presentation(PPTX_TEMPLATE_FILE)

  get_slide(
    prs,
    total_width_frac_list = total_width_frac,
    title = title,
    title_height_pt = title_height_pt,
    title_font_size_pt = title_font_size_pt,
    image_grid_list = image_grid_list,
    image_label_grid_list = label_grid_list,
    image_label_font_size_pt = image_label_font_size_pt,
    image_label_width_pt = image_label_width_pt,
    image_label_height_pt = image_label_height_pt,
    margin_top_font_size_pt = margin_top_font_size_pt,
    margin_left_font_size_pt = margin_left_font_size_pt,
    margin_top_height_pt = margin_top_height_pt,
    margin_left_width_pt = margin_left_width_pt,
    margin_left_spill_over_pt = margin_left_spill_over_pt[0],
    margin_right_spill_over_pt = margin_left_spill_over_pt[1],
    margin_label_left_list = margin_label_left_list,
    margin_label_top_list = margin_label_top_list,
    image_height_spacing_pt = image_height_spacing_pt,
    image_width_spacing_pt = image_width_spacing_pt,
    grid_height_spacing_pt = grid_height_spacing_pt,
  )
  
  log_utils.log(output)
  file_utils.make_parent_dir(output)
  prs.save(output)

if __name__ == '__main__':
  main(**parse_args())