import os
import sys
sys.path.append(os.path.abspath(os.path.join(os.path.dirname(__file__), '../utils/'))) # allow importing the utils dir
sys.path.append(os.path.abspath(os.path.join(os.path.dirname(__file__), '../2_graph_processing/'))) # allow importing the graphs dir

import numpy as np

import pptx
import pptx.util
import pptx.chart.data
import pptx.enum.shapes
import pptx.dml.color
import pptx.dml.effect

import pptx.shapes.connector
import pptx.enum.shapes
import pptx.enum.text
import pptx.enum.dml
import pptx.dml.color
import pptx.util

import PIL

import get_pptx_helpers

import library_constants

def get_legend_pptx(
  slide,
  title,
  items,
  x_pt,
  y_pt,
  stride_pt,
  title_width_pt,
  title_height_pt,
  title_x_offset_pt,
  item_width_pt,
  item_height_pt,
  label_width_pt,
  label_height_pt,
  legend_title_font_size_pt = 10,
  legend_label_font_size_pt = 8,
  orientation = 'v',
):
  get_pptx_helpers.add_textbox_pptx(
    slide,
    title,
    x_pt + title_x_offset_pt,
    y_pt,
    title_width_pt,
    title_height_pt,
    legend_title_font_size_pt,
    text_align = 'center',
  )
  y_pt += title_height_pt

  for _, item in enumerate(items):
    if item['type'] == 'circle':
      if orientation == 'v':
        shape = slide.shapes.add_shape(
          pptx.enum.shapes.MSO_SHAPE.OVAL,
          pptx.util.Pt(x_pt + item_width_pt / 2 - item['size'] / 2),
          pptx.util.Pt(y_pt + stride_pt / 2 - item['size'] / 2),
          pptx.util.Pt(item['size']),
          pptx.util.Pt(item['size']),
        )
      elif orientation == 'h':
        shape = slide.shapes.add_shape(
          pptx.enum.shapes.MSO_SHAPE.OVAL,
          pptx.util.Pt(x_pt + stride_pt / 2 - item['size'] / 2),
          pptx.util.Pt(y_pt + item_height_pt / 2 - item['size'] / 2),
          pptx.util.Pt(item['size']),
          pptx.util.Pt(item['size']),
        )
      else:
        raise Exception('Unknown orientation: ' + str(orientation))
      shape.fill.solid()
      shape.fill.fore_color.rgb = pptx.dml.color.RGBColor.from_string(
        item.get('color', 'FFFFFF').lstrip('#')
      )
      shape.line.color.rgb = pptx.dml.color.RGBColor.from_string(
        item.get('line_color', '000000').lstrip('#')
      )
      shape.line.width = pptx.util.Pt(item.get('line_width', 1))
    elif item['type'] == 'line':
      if orientation == 'v':
        shape = slide.shapes.add_connector(
          pptx.enum.shapes.MSO_CONNECTOR_TYPE.STRAIGHT,
          pptx.util.Pt(x_pt + item_width_pt / 2 - item['size'] / 2),
          pptx.util.Pt(y_pt + stride_pt / 2),
          pptx.util.Pt(x_pt + item_width_pt / 2 + item['size'] / 2),
          pptx.util.Pt(y_pt + stride_pt / 2),
        )
      elif orientation == 'h':
        shape = slide.shapes.add_connector(
          pptx.enum.shapes.MSO_CONNECTOR_TYPE.STRAIGHT,
          pptx.util.Pt(x_pt + stride_pt / 2 - item['size'] / 2),
          pptx.util.Pt(y_pt + item_height_pt / 2),
          pptx.util.Pt(x_pt + stride_pt / 2 + item['size'] / 2),
          pptx.util.Pt(y_pt + item_height_pt / 2),
        )
      else:
        raise Exception('Unknown orientation: ' + str(orientation))
      shape.line.color.rgb = pptx.dml.color.RGBColor.from_string(
        item.get('line_color', '000000')
      )
      shape.line.width = pptx.util.Pt(item.get('line_width', 1))
      shape.line.dash_style = {
        'dashed': pptx.enum.dml.MSO_LINE_DASH_STYLE.DASH,
        'solid': pptx.enum.dml.MSO_LINE_DASH_STYLE.SOLID,
      }[item.get('line_dash', 'solid')]
    else:
      raise ValueError('Unhandled item type: ' + str(item['type']))

    if orientation == 'v':
      shape = get_pptx_helpers.add_textbox_pptx(
        slide,
        item.get('text', ''),
        x_pt + item_width_pt,
        y_pt + stride_pt / 2 - label_height_pt / 2,
        label_width_pt,
        label_height_pt,
        legend_label_font_size_pt,
        text_align = 'left',
      )
      y_pt += stride_pt
    elif orientation == 'h':
      shape = get_pptx_helpers.add_textbox_pptx(
        slide,
        item.get('text', ''),
        x_pt + stride_pt / 2 - label_width_pt / 2,
        y_pt + item_height_pt,
        label_width_pt,
        label_height_pt,
        legend_label_font_size_pt,
        text_align = 'center',
      )
      x_pt += stride_pt
    else:
      raise Exception('Unknown orientation: ' + str(orientation))

  if orientation == 'h':
    y_pt += item_height_pt + label_height_pt
  return y_pt 

def get_edge_legend_pptx(
  slide,
  x_pt,
  y_pt,
  stride_pt,
  title_width_pt,
  title_height_pt,
  title_x_offset_pt,
  item_width_pt,
  item_height_pt,
  label_width_pt,
  label_height_pt,
  line_size_pt,
  line_width_pt,
  legend_title_font_size_pt = 10,
  legend_label_font_size_pt = 8,
  orientation = 'v',
):
  items = []

  # Substitution edges not being shown
  # items.append({
  #   'type': 'line',
  #   'text': 'Substitution',
  #   'color': '000000',
  #   'line_dash': 'solid',
  #   'line_width': line_width_pt,
  #   'size': line_size_pt,
  # })
  items.append({
    'type': 'line',
    'text': '1 nt. in/del',
    'color': '000000',
    'line_dash': 'solid',
    # In/del edges being shown solid
    # 'line_dash': 'dash',
    'line_width': line_width_pt,
    'size': line_size_pt,
  })
  return get_legend_pptx(
    slide = slide,
    title = 'Edges',
    items = items,
    x_pt = x_pt,
    y_pt = y_pt,
    stride_pt = stride_pt,
    title_width_pt = title_width_pt,
    title_height_pt = title_height_pt,
    title_x_offset_pt = title_x_offset_pt,
    item_width_pt = item_width_pt,
    item_height_pt = item_height_pt,
    label_width_pt = label_width_pt,
    label_height_pt = label_height_pt,
    legend_title_font_size_pt = legend_title_font_size_pt,
    legend_label_font_size_pt = legend_label_font_size_pt,
    orientation = orientation,
  )

def get_variation_color_legend_pptx(
  slide,
  x_pt,
  y_pt,
  stride_pt,
  title_width_pt,
  title_height_pt,
  title_x_offset_pt,
  item_width_pt,
  item_height_pt,
  label_width_pt,
  label_height_pt,
  variation_types,
  node_size_pt,
  line_width_pt,
  legend_title_font_size_pt = 10,
  legend_label_font_size_pt = 8,
  orientation = 'v',
):
  items = []
  for var_type in variation_types:
    items.append({
      'type': 'circle',
      'size': node_size_pt,
      'text': library_constants.VARIATION_TYPES[var_type]['label'],
      'color': library_constants.VARIATION_TYPES[var_type]['color'],
      'line_width': line_width_pt,
    })
  return get_legend_pptx(
    slide = slide,
    title = 'Variation Types',
    items = items,
    stride_pt = stride_pt,
    x_pt = x_pt,
    y_pt = y_pt,
    title_width_pt = title_width_pt,
    title_height_pt = title_height_pt,
    title_x_offset_pt = title_x_offset_pt,
    item_width_pt = item_width_pt,
    item_height_pt = item_height_pt,
    label_width_pt = label_width_pt,
    label_height_pt = label_height_pt,
    legend_title_font_size_pt = legend_title_font_size_pt,
    legend_label_font_size_pt = legend_label_font_size_pt,
    orientation = orientation,
  )

def get_outline_legend_pptx(
  slide,
  x_pt,
  y_pt,
  stride_pt,
  title_width_pt,
  title_height_pt,
  title_x_offset_pt,
  item_width_pt,
  item_height_pt,
  label_width_pt,
  label_height_pt,
  node_size_pt,
  line_width_pt,
  legend_title_font_size_pt = 10,
  legend_label_font_size_pt = 8,
  orientation = 'v',
):
  items = []
  items.append({
    'type': 'circle',
    'size': node_size_pt,
    'text': 'Reference',
    'color': 'FFFFFF',
    'line_color': library_constants.REFERENCE_OUTLINE_COLOR,
    'line_width': 2 * line_width_pt,
  })
  items.append({
    'type': 'circle',
    'size': node_size_pt,
    'text': 'Non-reference',
    'color': 'FFFFFF',
    'line_color': library_constants.DEFAULT_OUTLINE_COLOR,
    'line_width': line_width_pt,
  })
  return get_legend_pptx(
    slide = slide,
    title = 'Outline',
    items = items,
    x_pt = x_pt,
    y_pt = y_pt,
    stride_pt = stride_pt,
    title_width_pt = title_width_pt,
    title_height_pt = title_height_pt,
    title_x_offset_pt = title_x_offset_pt,
    item_width_pt = item_width_pt,
    item_height_pt = item_height_pt,
    label_width_pt = label_width_pt,
    label_height_pt = label_height_pt,
    legend_title_font_size_pt = legend_title_font_size_pt,
    legend_label_font_size_pt = legend_label_font_size_pt,
    orientation = orientation,
  )

def get_node_legend_pptx(
  slide,
  x_pt,
  y_pt,
  stride_pt,
  title_width_pt,
  title_height_pt,
  title_x_offset_pt,
  item_width_pt,
  item_height_pt,
  label_width_pt,
  label_height_pt,
  node_size_pt,
  line_width_pt,
  legend_title_font_size_pt = 10,
  legend_label_font_size_pt = 8,
  orientation = 'v',
):
  items = [
    {
      'type': 'circle',
      'size': node_size_pt,
      'text': library_constants.VARIATION_TYPES['insertion']['label'],
      'color': library_constants.VARIATION_TYPES['insertion']['color'],
      'line_width': line_width_pt,
    },
    {
      'type': 'circle',
      'size': node_size_pt,
      'text': library_constants.VARIATION_TYPES['deletion']['label'],
      'color': library_constants.VARIATION_TYPES['deletion']['color'],
      'line_width': line_width_pt,
    },
    {
      'type': 'circle',
      'size': node_size_pt,
      'text': library_constants.REFERENCE_DESCRIPTION,
      'color': library_constants.DEFAULT_NODE_COLOR,
      'line_width': 2 * line_width_pt,
      'line_color': library_constants.REFERENCE_OUTLINE_COLOR,
    },
  ]
  return get_legend_pptx(
    slide = slide,
    title = 'Vertices',
    items = items,
    stride_pt = stride_pt,
    x_pt = x_pt,
    y_pt = y_pt,
    title_width_pt = title_width_pt,
    title_height_pt = title_height_pt,
    title_x_offset_pt = title_x_offset_pt,
    item_width_pt = item_width_pt,
    item_height_pt = item_height_pt,
    label_width_pt = label_width_pt,
    label_height_pt = label_height_pt,
    legend_title_font_size_pt = legend_title_font_size_pt,
    legend_label_font_size_pt = legend_label_font_size_pt,
    orientation = orientation,
  )

def get_size_legend_pptx(
  slide,
  x_pt,
  y_pt,
  stride_pt,
  title_width_pt,
  title_height_pt,
  title_x_offset_pt,
  item_width_pt,
  item_height_pt,
  label_width_pt,
  label_height_pt,
  node_size_min_freq,
  node_size_max_freq,
  node_size_min_pt,
  node_size_max_pt,
  line_width_pt,
  legend_title_font_size_pt = 10,
  legend_label_font_size_pt = 8,
  orientation = 'v',
):
  node_size_min_freq_log10 = int(np.round(np.log10(node_size_min_freq)))
  node_size_max_freq_log10 = int(np.round(np.log10(node_size_max_freq)))

  num_legend_items = node_size_max_freq_log10 - node_size_min_freq_log10 + 1

  items = []
  for i in range(num_legend_items):
    freq_log10 = node_size_min_freq_log10 + i
    if num_legend_items == 1:
      size = node_size_min_pt
    else:
      size = node_size_min_pt + (
        i * (node_size_max_pt - node_size_min_pt) /
        (num_legend_items - 1)
      )
    if freq_log10 == 0:
      text = '1'
    else:
      text = f'10' + {-1: '⁻¹', -2: '⁻²', -3: '⁻³', -4: '⁻⁴', -5: '⁻⁵', -6: '⁻⁶'}[freq_log10]
    if i == 0:
      text = '≤' + text
    items.append({
      'type': 'circle',
      'size': size,
      'text': text,
      'color': 'FFFFFF',
      'line_width': line_width_pt,
    })
  items = items[::-1] # Show largest to smallest
  return get_legend_pptx(
    slide = slide,
    title = 'Sequence frequency',
    items = items,
    x_pt = x_pt,
    y_pt = y_pt,
    stride_pt = stride_pt,
    title_width_pt = title_width_pt,
    title_height_pt = title_height_pt,
    title_x_offset_pt = title_x_offset_pt,
    item_width_pt = item_width_pt,
    item_height_pt = item_height_pt,
    label_width_pt = label_width_pt,
    label_height_pt = label_height_pt,
    legend_title_font_size_pt = legend_title_font_size_pt,
    legend_label_font_size_pt = legend_label_font_size_pt,
    orientation = orientation,
  )

def get_freq_ratio_legend_pptx(
  slide,
  x_pt,
  y_pt,
  title,
  title_width_pt,
  title_height_pt,
  title_x_offset_pt,
  label_width_pt,
  label_height_pt,
  color_bar_minor_axis_pt,
  color_bar_major_axis_pt,
  color_bar_file,
  legend_title_font_size_pt = 8,
  legend_label_font_size_pt = 10,
  orientation = 'v',
  fixed_aspect = False,
):

  get_pptx_helpers.add_textbox_pptx(
    slide,
    title,
    x_pt + title_x_offset_pt,
    y_pt,
    title_width_pt,
    title_height_pt,
    legend_title_font_size_pt,
  )
  y_pt += title_height_pt

  if fixed_aspect:
    image = PIL.Image.open(color_bar_file)
    image_width_px = image.width
    image_height_px = image.height
    ratio = min(
      color_bar_width_pt / image_width_px,
      color_bar_height_pt / image_height_px,
    )
    color_bar_width_pt = ratio * image_width_px
    color_bar_height_pt = ratio * image_height_px

  if orientation == 'v':
    color_bar = slide.shapes.add_picture(
      color_bar_file,
      pptx.util.Pt(x_pt),
      pptx.util.Pt(y_pt),
      pptx.util.Pt(color_bar_minor_axis_pt),
      pptx.util.Pt(color_bar_major_axis_pt),
    )
  elif orientation == 'h':
    color_bar = slide.shapes.add_picture(
      color_bar_file,
      pptx.util.Pt(x_pt + color_bar_major_axis_pt / 2 - color_bar_minor_axis_pt / 2),
      pptx.util.Pt(y_pt - color_bar_major_axis_pt / 2),
      pptx.util.Pt(color_bar_minor_axis_pt),
      pptx.util.Pt(color_bar_major_axis_pt),
    )
    color_bar.rotation = -90
  else:
    raise Exception('Unknown orientation: ' + str(orientation))
  num_ticks = len(library_constants.FREQ_RATIO_COLOR_BAR_TICK_TEXT)
  for i, text in enumerate(library_constants.FREQ_RATIO_COLOR_BAR_TICK_TEXT):
    i_rev = len(library_constants.FREQ_RATIO_COLOR_BAR_TICK_TEXT) - i - 1 # descending order
    if orientation == 'v':
      x_label_pt = x_pt + color_bar_minor_axis_pt / 2
      y_label_pt = (
        y_pt + (i_rev / (num_ticks - 1)) * color_bar_major_axis_pt - label_height_pt / 2
      )
      text_align = 'left'
    elif orientation == 'h':
      x_label_pt = (
        x_pt + (i_rev / (num_ticks - 1)) * color_bar_major_axis_pt - label_width_pt / 2
      )
      y_label_pt = y_pt
      text_align = 'center'
    else:
      raise Exception('Unknown orientation: ' + str(orientation))
    get_pptx_helpers.add_textbox_pptx(
      slide = slide,
      text = text,
      x_pt = x_label_pt,
      y_pt = y_label_pt,
      width_pt = label_width_pt,
      height_pt = label_height_pt,
      font_size_pt = legend_label_font_size_pt,
      text_align = text_align,
    )
  y_pt += color_bar_major_axis_pt

  return y_pt

